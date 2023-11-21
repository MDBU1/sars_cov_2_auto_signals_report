import os
import re
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from datetime import date
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pd2ppt import df_to_table

from modules.class_signal_stats import ClassSignalStatistics


# %%


class ClassPresentationSetup(ClassSignalStatistics):
    def __init__(self, class_signal):
        super().__init__(class_signal)
        self.str_prs_signal_call = self.meth_return_prs_name()
        self.str_prs_signal_call_save = self.meth_return_prs_save_name()
        self.meth_generate_slide_titles()
        self.meth_generate_slide_epi()
        self.meth_generate_slide_genetic()
        self.meth_generate_slide_overview()
        # self.meth_generate_slide_deck()

    # Slide Components
    def meth_return_prs_name(self):
        mutations = " ".join(self.mutations)
        if mutations is "no_specified_mutations":
            str_mutations = " "
        else:
            str_mutations = mutations.replace("_", " ")
        if self.lineage is "no_specified_lineage":
            str_lineage = " "
        else:
            str_lineage = self.lineage.replace("_", " ")

        temp_str_signal_call = f"{self.signal_no} {str_lineage} {str_mutations}"
        return temp_str_signal_call

    def meth_return_prs_save_name(self):
        str_replacements = {' ': '_', ':': '.'}
        pattern = '|'.join(str_replacements.keys())
        str_replaced_signal = re.sub(pattern, lambda match: str_replacements[match.group(0)], self.str_prs_signal_call)
        return str_replaced_signal

    def meth_slide_header(self, input_slide):
        slide1_header = input_slide.shapes.add_textbox(Cm(1), Cm(0), Cm(10), Cm(0.1))
        tf = slide1_header.text_frame
        p = tf.paragraphs[0]
        p.text = f"OFFICIAL SENSITIVE " + "{:%Y%m%d}".format(date.today()) + f" {self.meth_return_prs_name()}"
        p.font.size = Pt(10)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(0, 124, 145)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    def meth_slide_footer(self, input_slide):
        slide1_footer = input_slide.shapes.add_textbox(Cm(6), Cm(18), Cm(10), Cm(0.1))
        tf = slide1_footer.text_frame
        p = tf.paragraphs[0]
        p.text = f"OFFICIAL SENSITIVE SARS-CoV2 Horizon Scanning" + " {:%Y%m%d}".format(date.today())
        # + f" {self.meth_return_prs_name()}"
        p.font.size = Pt(14)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(0, 124, 145)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    def meth_generate_slide(self):
        prs = Presentation(self.prs_slide_main)
        slide1_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(slide1_layout)
        slide1 = prs.slides[0]
        self.meth_slide_header(slide1)
        self.meth_slide_footer(slide1)
        # self.meth_add_summary_box(slide1)
        del prs.slides._sldIdLst[1]
        return prs

    def meth_add_slide_title(self, input_slide, slide_name):
        slide1_title = input_slide.shapes.add_textbox(Cm(4.5), Cm(0.45), Cm(22), Cm(5))
        tf = slide1_title.text_frame
        p = tf.paragraphs[0]
        p.text = f"SARs-CoV2 SIGNAL ANALYSIS: {self.data_region.upper()} {slide_name}"
        p.font.size = Pt(32)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(0, 124, 145)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        return 0

    # Title Slide
    def meth_generate_slide_titles(self):
        prs = Presentation(self.prs_slide_title)

        slide1_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(slide1_layout)

        slide1_title = slide1.shapes.add_textbox(Cm(4), Cm(6.5), Cm(15), Cm(7.5))
        tf = slide1_title.text_frame
        p = tf.paragraphs[0]
        p.text = "UKHSA \nSARs-CoV2 Horizon Scanning"
        p.font.size = Pt(36)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(0, 124, 145)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        sub = tf.add_paragraph()
        sub.text = f"\nSignal Analysis: {self.str_prs_signal_call}"
        sub.font.size = Pt(36)
        sub.font.name = 'Arial'
        sub.font.color.rgb = RGBColor(0, 124, 145)

        sub_date = tf.add_paragraph()
        sub_date.text = "{:%Y%m%d}".format(date.today())
        sub_date.font.size = Pt(11)
        sub_date.font.name = 'Arial'
        sub_date.font.color.rgb = RGBColor(0, 124, 145)

        del prs.slides._sldIdLst[0]  # deletes base "template" slide
        prs.save(self.path_save_dir + "/" + "{:%Y%m%d}".format(date.today()) + "_SA_presentation_title_" +
                 self.str_prs_signal_call_save + ".pptx")
        return prs

    # Overview Slide
    def meth_add_signal_vs_top_lineages(self, input_slide):
        file = self.path_save_dir + "/fig_lineage_percentages.png"
        if not os.path.isfile(file):
            print("lineage percentages figure not found")
        else:
            left = Cm(1.5)
            top = Cm(2.5)
            width = Cm(15)
            height = Cm(9)
            pic = input_slide.shapes.add_picture(file, left, top, width, height)

    def meth_add_signal_weekly_cumulative_fig(self, input_slide):
        file = self.path_save_dir + "/fig_weekly_cumulative_epi.png"
        if not os.path.isfile(file):
            print("cumulative epi curve not found")
        else:
            left = Cm(17)
            top = Cm(2.5)
            width = Cm(15)
            height = Cm(9)
            pic = input_slide.shapes.add_picture(file, left, top, width, height)

    def meth_iter_cells(self, table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    def meth_insert_signal_update_table(self, slide):
        df = self.table_update_summary
        x, y, cx, cy = Cm(1.5), Cm(11.5), Cm(31), Cm(4)
        # tbl1 = df_to_table(slide, dt, left, top, width, height, name='tbl1')
        df_to_table(slide, df, x, y, cx, cy, name='table_signal')
        for shape in slide.shapes:
            # print(shape.shape_type)
            if shape.has_table:
                table = shape.table
                table.columns[0].width = Cm(3.5)
                table.columns[1].width = Cm(3.5)
                table.columns[2].width = Cm(24)
                for cell in self.meth_iter_cells(table):
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(13)

    def meth_generate_slide_overview(self):
        prs = self.meth_generate_slide()
        slide = prs.slides[0]
        self.meth_add_signal_vs_top_lineages(slide)
        self.meth_add_signal_weekly_cumulative_fig(slide)
        self.meth_insert_signal_update_table(slide)
        self.meth_add_slide_title(slide, "Overview        ")
        prs.save(self.path_save_dir + "/" + "{:%Y%m%d}".format(date.today()) + "_SA_overview_" +
                 self.str_prs_signal_call_save + ".pptx")
        return prs

    # EPI SLide
    def meth_add_epi_curve(self, input_slide):
        file = self.path_save_dir + "/fig_epi_curve.png"
        if not os.path.isfile(file):
            print("epi. curve not found")
        else:
            left = Cm(11.8)
            top = Cm(0.25)
            width = Cm(22)
            height = Cm(9)
            pic = input_slide.shapes.add_picture(file, left, top, width, height)

    def meth_add_prop_epi_curve(self, input_slide):
        file_prop_region_epi_curve = self.path_save_dir + "/fig_percentage_region_weekly_epi_curves.png"
        if not os.path.isfile(file_prop_region_epi_curve):
            print("prop. epi. curves not found")
        else:
            left = Cm(11.8)
            top = Cm(9.25)
            width = Cm(22)
            height = Cm(9)
            pic = input_slide.shapes.add_picture(file_prop_region_epi_curve, left, top, width, height)

    def meth_add_map_int(self, input_slide, input_file):
        # file = self.path_save_dir + "/fig_int_map.png"
        file = input_file
        if not os.path.isfile(file):
            print("int map not found")
        else:
            left = Cm(1.5)
            top = Cm(9.25)
            width = Cm(10.5)
            height = Cm(9)
            pic = input_slide.shapes.add_picture(file, left, top, width, height)

    def meth_add_epi_slide_summary_box(self, input_slide):
        # summary text
        rect0 = input_slide.shapes.add_shape(  # Add Shape object➀
            MSO_SHAPE.RECTANGLE,  # Specify the shape type as "Rounded Rectangle"
            Cm(1.5), Cm(3.5),  # Insertion position is specified as the upper left coordinate of the shape
            Cm(10.25), Cm(6))  # Specify the width and height of the inserted shape
        fill = rect0.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        line = rect0.line
        line.color.rgb = RGBColor(0, 124, 145)
        # self.func_summary_epi_text(rect0)
        return rect0

    def meth_return_epi_caption_text(self, input_slide):
        caption = input_slide.shapes.add_textbox(Cm(20.5), Cm(18.25), Cm(10), Cm(0.1))
        tf = caption.text_frame
        p = tf.paragraphs[0]
        p.text = f"Numbers included next to the region names above each plot represent the total number of sequences " \
                 f"designated as matching signal \n specification across the analysis period (plot A) and the total " \
                 f"number of sequences for each region across the analysis period (plot B)"
        p.font.size = Pt(7)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(0, 0, 0)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    def meth_return_epi_slide_summary_text(self, insert_text_box):
        insert_text_box.text = (f"Total sequences: {self.stat_total_sequences}"
                                + "\n" +
                                f"Highest reporting region: {self.stat_highest_region}"
                                + "\n" +
                                f"Highest reporting country: {self.stat_highest_country}"
                                + "\n" +
                                f"Highest weekly sequences: {self.stat_highest_weekly_sequences}"
                                + "\n" +
                                # f"Highest weekly percentage per region: {self.stat_highest_week_percent_per_region}"
                                # + "\n" +
                                f"Significant contributors: {self.stat_subregion_outliers}"
                                + "\n" +
                                f"Number of sub-lineages: {len(self.df_lineage_subs) - 1}")

        for x in range(len(insert_text_box.text_frame.paragraphs)):
            p = insert_text_box.text_frame.paragraphs[x]
            p.font.size = Pt(13.5)
            p.font.name = 'Arial'
            p.font.color.rgb = RGBColor(0, 0, 0)

    def meth_generate_slide_epi(self):
        prs = self.meth_generate_slide()
        slide = prs.slides[0]
        if self.data_region == "eng":
            file = self.path_save_dir + "/map_UTLA_counts.png"
            self.meth_add_map_int(slide, file)
        elif self.data_region == "int":
            file = self.path_save_dir + "/fig_int_map.png"
            self.meth_add_map_int(slide, file)
        self.meth_add_epi_curve(slide)
        self.meth_add_prop_epi_curve(slide)
        self.meth_add_slide_title(slide, "\nEpidemiology                                                              "
                                         "     ")
        rect = self.meth_add_epi_slide_summary_box(slide)
        self.meth_return_epi_slide_summary_text(rect)
        self.meth_return_epi_caption_text(slide)
        prs.save(self.path_save_dir + "/" + "{:%Y%m%d}".format(date.today()) + "_SA_presentation_epi_" +
                 self.str_prs_signal_call_save + ".pptx")
        return prs

    # Genetic SLide

    def meth_add_mutation_profile_table(self, input_slide):
        left = Cm(1.5)
        top = Cm(2.25)
        width = Cm(30)
        height = Cm(9)
        table = df_to_table(input_slide, self.df_profile_mutation_ordered, left, top, width, height)
        for shape in input_slide.shapes:
            if shape.has_table:
                table = shape.table
                table.columns[0].width = Cm(2)
                table.columns[1].width = Cm(30)

                def iter_cells(table):
                    for row in table.rows:
                        for cell in row.cells:
                            yield cell

                for cell in iter_cells(table):
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(12)

    def meth_add_genetics_slide_summary_box(self, input_slide):
        rect0 = input_slide.shapes.add_shape(  # Add Shape object➀
            MSO_SHAPE.RECTANGLE,  # Specify the shape type as "Rounded Rectangle"
            Cm(19), Cm(17),  # Insertion position is specified as the upper left coordinate of the shape
            Cm(14.5), Cm(2))  # Specify the width and height of the inserted shape
        fill = rect0.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        line = rect0.line
        line.color.rgb = RGBColor(0, 124, 145)
        # self.func_summary_epi_text(rect0)
        return rect0

    def meth_return_genetics_slide_summary_text(self, insert_text_box):
        insert_text_box.text = (f"Unique mutation/s compared with parent lineage: "
                                + "\n" +
                                f"Unique mutation/s compared with top ten circulating lineages: "
                                f"{self.stat_no_unique_mutations}, {self.stat_unique_mutations}"
                                + "\n" +
                                f"Mutations conserved in <=0.5 top circulating lineages: "
                                f"{self.stat_no_uncommon_mutations}, {self.stat_uncommon_mutations}"
                                + "\n" +
                                f"Note. Mutations conserved based within analysis time-period.")

        for x in range(len(insert_text_box.text_frame.paragraphs)):
            p = insert_text_box.text_frame.paragraphs[x]
            p.font.size = Pt(10)
            p.font.name = 'Arial'
            p.font.color.rgb = RGBColor(0, 0, 0)

    def meth_generate_slide_genetic(self):
        prs = self.meth_generate_slide()
        slide = prs.slides[0]
        self.meth_add_mutation_profile_table(slide)
        rect = self.meth_add_genetics_slide_summary_box(slide)
        self.meth_return_genetics_slide_summary_text(rect)
        self.meth_add_slide_title(slide, "Mutation Profile")
        prs.save(self.path_save_dir + "/" + "{:%Y%m%d}".format(date.today()) + "_SA_presentation_genetic_" +
                 self.str_prs_signal_call_save + ".pptx")
        return prs

    def meth_generate_slide_deck(self):
        prs1 = self.meth_generate_slide_titles()
        prs2 = self.meth_generate_slide_overview()
        prs3 = self.meth_generate_slide_epi()
        prs4 = self.meth_generate_slide_genetic()

        prs4.save(self.path_save_dir + "/" + "{:%Y%m%d}".format(date.today()) + "_SA_presentation_" +
                  self.str_prs_signal_call_save + ".pptx")
